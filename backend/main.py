"""
주일예배 PPT 생성기 - FastAPI 백엔드
"""

import os
import re
import shutil
import tempfile
from pathlib import Path

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation

from worship_core import (
    _replace_date_in_prs,
    find_hymn_blocks,
    find_responsory_block,
    replace_slides_zip,
    parse_godpia_text,
    apply_bible_text,
    replace_sermon_title,
    apply_sermon_song,
    parse_choir_input,
    apply_choir_lyrics,
    parse_praise_input,
    apply_praise_lyrics,
)
from drive_client import download_hymn, download_responsory, download_template

app = FastAPI(title="주일예배 PPT 생성기")

# 프론트엔드 정적 파일 서빙
_static_dir = os.path.join(os.path.dirname(__file__), 'static')
if os.path.exists(_static_dir):
    app.mount('/static', StaticFiles(directory=_static_dir), name='static')

@app.get('/', response_class=HTMLResponse)
def root():
    idx = os.path.join(_static_dir, 'index.html')
    return open(idx, encoding='utf-8').read()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def _cleanup(path: str):
    """임시 디렉토리 정리"""
    try:
        shutil.rmtree(path, ignore_errors=True)
    except Exception:
        pass


@app.get("/health")
def health():
    return {"status": "ok"}


@app.post("/generate")
async def generate_ppt(
    background_tasks: BackgroundTasks,
    part: str = Form(...),                          # '1bu' or '2bu'
    date_str: str = Form(""),
    hymn1: str = Form(""),
    resp_num: str = Form(""),
    hymn2: str = Form(""),
    sermon_title: str = Form(""),
    sermon_ref: str = Form(""),
    bible_label: str = Form(""),
    bible_text: str = Form(""),
    sermon_song_title: str = Form(""),
    choir_lyrics_1bu: str = Form(""),
    choir_lyrics_2bu: str = Form(""),
    choir_name_1bu: str = Form("호산나찬양대"),
    choir_name_2bu: str = Form("시온찬양대"),
    praise_text: str = Form(""),
    template_file_id: str = Form(""),              # 사용자 지정 템플릿 Drive ID
    sermon_song_file: UploadFile = File(None),     # 설교찬양 .ppt/.pptx 업로드
):
    if part not in ('1bu', '2bu'):
        raise HTTPException(status_code=400, detail="part는 '1bu' 또는 '2bu'여야 합니다.")

    tmpdir = tempfile.mkdtemp(prefix='worship_')
    background_tasks.add_task(_cleanup, tmpdir)

    log = []

    def _log(msg: str):
        log.append(msg)

    try:
        # ── 템플릿 다운로드 ───────────────────────────
        _log(f"📥 템플릿 다운로드 중...")
        tmpl_path = download_template(part, tmpdir, custom_file_id=template_file_id or None)
        if not tmpl_path:
            raise HTTPException(status_code=500, detail="템플릿 파일을 Drive에서 찾을 수 없습니다.")
        _log(f"✅ 템플릿: {Path(tmpl_path).name}")

        # 출력 파일 경로
        date_clean = re.sub(r"[^\d]", "", date_str)[:8] or "00000000"
        label = "1부" if part == "1bu" else "2부"
        out_filename = f"{date_clean}-주일예배({label}).pptx"
        out_path = str(Path(tmpdir) / out_filename)
        shutil.copy2(tmpl_path, out_path)

        # ── 날짜 치환 ────────────────────────────────
        if date_str:
            _log(f"📅 날짜: {date_str}")
            prs = Presentation(out_path)
            _replace_date_in_prs(prs, date_str)
            prs.save(out_path)

        # ── 찬송가① 교체 ─────────────────────────────
        if hymn1.strip():
            num = int(hymn1.strip())
            _log(f"📥 찬송가① {num}장 다운로드 중...")
            hymn1_path = download_hymn(num, tmpdir)
            if hymn1_path:
                _log(f"🎵 찬송가① {num}장 교체 중...")
                prs = Presentation(out_path)
                blocks = find_hymn_blocks(prs)
                if blocks:
                    s, c = blocks[0]
                    replace_slides_zip(out_path, s, c, hymn1_path, out_path)
                    _log(f"   → 슬라이드 {s+1}~{s+c} 교체 완료")
                else:
                    _log("   ⚠️ 찬송가① 블록을 찾지 못했습니다.")
            else:
                _log(f"⚠️ 찬송가① {num}장 파일을 Drive에서 찾지 못했습니다.")

        # ── 교독문 교체 ───────────────────────────────
        if resp_num.strip():
            num = int(resp_num.strip())
            _log(f"📥 교독문 {num}번 다운로드 중...")
            resp_path = download_responsory(num, tmpdir)
            if resp_path:
                _log(f"📋 교독문 {num}번 교체 중...")
                prs = Presentation(out_path)
                block = find_responsory_block(prs)
                if block:
                    s, c = block
                    replace_slides_zip(out_path, s, c, resp_path, out_path)
                    _log(f"   → 슬라이드 {s+1}~{s+c} 교체 완료")
                else:
                    _log("   ⚠️ 교독문 블록을 찾지 못했습니다.")
            else:
                _log(f"⚠️ 교독문 {num}번 파일을 Drive에서 찾지 못했습니다.")

        # ── 찬송가② 교체 ─────────────────────────────
        if hymn2.strip():
            num = int(hymn2.strip())
            _log(f"📥 찬송가② {num}장 다운로드 중...")
            hymn2_path = download_hymn(num, tmpdir)
            if hymn2_path:
                _log(f"🎵 찬송가② {num}장 교체 중...")
                prs = Presentation(out_path)
                blocks = find_hymn_blocks(prs)
                if len(blocks) >= 2:
                    s, c = blocks[1]
                    replace_slides_zip(out_path, s, c, hymn2_path, out_path)
                    _log(f"   → 슬라이드 {s+1}~{s+c} 교체 완료")
                else:
                    _log("   ⚠️ 찬송가② 블록을 찾지 못했습니다.")
            else:
                _log(f"⚠️ 찬송가② {num}장 파일을 Drive에서 찾지 못했습니다.")

        # ── 성경 봉독 ─────────────────────────────────
        if bible_text.strip():
            if bible_label.strip():
                book_chapter = bible_label.strip()
            else:
                m = re.match(r'(.+?)\s+(\d+):\d+', sermon_ref.strip() if sermon_ref else "")
                book_chapter = f"{m.group(1)} {m.group(2)}" if m else sermon_ref.strip()
            _log(f"📜 성경봉독: {book_chapter}")
            verses = parse_godpia_text(bible_text.strip(), book_chapter)
            prs = Presentation(out_path)
            applied = apply_bible_text(prs, verses)
            prs.save(out_path)
            _log(f"   → {applied}절 입력 완료")

        # ── 설교 제목 ─────────────────────────────────
        if sermon_title.strip():
            _log(f"📖 설교 제목: {sermon_title}")
            prs = Presentation(out_path)
            ok = replace_sermon_title(prs, sermon_title.strip(), sermon_ref.strip())
            prs.save(out_path)
            if not ok:
                _log("   ⚠️ 설교 제목 슬라이드를 찾지 못했습니다.")

        # ── 설교찬양 악보 ─────────────────────────────
        if sermon_song_file and sermon_song_file.filename:
            _log(f"🎼 설교찬양: {sermon_song_title or '(제목없음)'}")
            suffix = Path(sermon_song_file.filename).suffix.lower()
            score_path = str(Path(tmpdir) / f'sermon_song{suffix}')
            with open(score_path, 'wb') as f:
                f.write(await sermon_song_file.read())
            prs = Presentation(out_path)
            applied_ss = apply_sermon_song(prs, score_path, sermon_song_title.strip())
            if applied_ss == -1:
                _log("   ⚠️ .ppt 변환 실패. LibreOffice가 설치되어 있는지 확인하세요.")
            elif applied_ss:
                prs.save(out_path)
                _log(f"   → {applied_ss}슬라이드 삽입 완료")
            else:
                _log("   ⚠️ 설교찬양 슬라이드를 찾지 못했습니다.")

        # ── 성가대 자막 ───────────────────────────────
        choir_lyrics = choir_lyrics_1bu if part == '1bu' else choir_lyrics_2bu
        choir_name   = choir_name_1bu   if part == '1bu' else choir_name_2bu
        if choir_lyrics.strip():
            lines = choir_lyrics.strip().splitlines()
            first_line = lines[0].strip() if lines else ""
            if first_line.startswith('[') and first_line.endswith(']'):
                choir_song_title = first_line[1:-1]
                remaining = '\n'.join(lines[1:]).strip()
                choir_blocks = parse_choir_input(remaining)
            else:
                choir_song_title = ""
                choir_blocks = parse_choir_input(choir_lyrics)
            _log(f"🎵 성가대 자막: {len(choir_blocks)}슬라이드")
            prs = Presentation(out_path)
            applied = apply_choir_lyrics(prs, choir_blocks, choir_name, choir_song_title)
            prs.save(out_path)
            _log(f"   → {applied}슬라이드 입력 완료")

        # ── 찬양팀 가사 (2부 전용) ────────────────────
        if part == '2bu' and praise_text.strip():
            songs = parse_praise_input(praise_text.strip())
            # \x00(임시 줄바꿈 치환자) → \n 복원
            for song in songs:
                song['title'] = song['title'].replace('\x00', '\n')
                song['slides'] = [
                    [line.replace('\x00', '\n') for line in slide]
                    for slide in song['slides']
                ]
            _log(f"✨ 찬양팀 가사: {len(songs)}곡")
            applied = apply_praise_lyrics(out_path, songs)
            _log(f"   → {applied}슬라이드 입력 완료")

        _log(f"\n✅ {label} 생성 완료!")

        # pptx 파일을 base64로 인코딩해서 JSON으로 반환
        import base64
        with open(out_path, 'rb') as f:
            file_b64 = base64.b64encode(f.read()).decode('utf-8')

        from fastapi.responses import JSONResponse
        return JSONResponse({
            "filename": out_filename,
            "file": file_b64,
            "log": log
        })

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        detail = f"{e}\n{traceback.format_exc()}"
        raise HTTPException(status_code=500, detail=detail)


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
