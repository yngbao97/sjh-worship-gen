"""
주일예배 PPT 자동 생성 - 핵심 로직 (웹 버전)
원본 worship_ppt.py에서 tkinter/GUI 제거, FastAPI 백엔드용
"""

import os, sys, json, re, shutil, zipfile, threading, datetime, tempfile, copy
from pathlib import Path
from lxml import etree
from pptx.oxml import parse_xml

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
# ───────────────────────────────────────────────
# 네임스페이스 상수
# ───────────────────────────────────────────────
NS_P   = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_REL = 'http://schemas.openxmlformats.org/package/2006/relationships'
NS_CT  = 'http://schemas.openxmlformats.org/package/2006/content-types'



# ───────────────────────────────────────────────
# 파일 탐색 유틸
# ───────────────────────────────────────────────

def get_hymn_subfolder_name(number: int) -> str:
    """찬송가 번호 → 서브폴더명. '1~100장' 과 '1~100' 둘 다 시도."""
    for start, end in HYMN_RANGES:
        if start <= number <= end:
            return f"{start}~{end}"
    return ""

def _find_hymn_subfolder(root_folder: str, number: int) -> str | None:
    """실제 존재하는 서브폴더를 탐색 (이름 패턴이 달라도 대응)."""
    from pathlib import Path
    root = Path(root_folder)
    if not root.exists():
        return None
    candidates = [
        f"{get_hymn_subfolder_name(number)}장",   # '1~100장'
        get_hymn_subfolder_name(number),           # '1~100'
    ]
    for name in candidates:
        if name and (root / name).is_dir():
            return str(root / name)
    # 위 후보가 없으면 폴더 내 모든 디렉터리를 순회해서 번호 범위 포함 여부 확인
    for d in root.iterdir():
        if not d.is_dir(): continue
        m = __import__('re').search(r'(\d+)\D+(\d+)', d.name)
        if m:
            s, e = int(m.group(1)), int(m.group(2))
            if s <= number <= e:
                return str(d)
    return None

def find_hymn_path(root_folder: str, number: str) -> str | None:
    """찬송가 번호로 PPT 파일 경로 반환. 서브폴더명/파일명 패턴 자동 탐색."""
    try:
        num = int(number.strip())
    except ValueError:
        return None
    subfolder = _find_hymn_subfolder(root_folder, num)
    if not subfolder:
        return None
    # 파일명 후보 순서대로 시도
    for fname in [f"{num}장.pptx", f"{num}.pptx", f"찬송가{num}장.pptx", f"찬송가_{num}장.pptx"]:
        p = Path(subfolder) / fname
        if p.exists():
            return str(p)
    return None

def find_responsory_path(folder: str, number: str) -> str | None:
    """교독문 번호로 PPT 파일 경로 반환. 파일명 패턴 자동 탐색."""
    try:
        num = int(number.strip())
    except ValueError:
        return None
    folder_path = Path(folder)
    # 파일명 후보 순서대로 시도
    for fname in [
        f"교독문 {num}번.pptx",    # 공백
        f"교독문_{num}번.pptx",    # 언더바
        f"교독문{num}번.pptx",     # 붙여쓰기
        f"교독문 {num:03d}번.pptx", # 공백+3자리
        f"교독문_{num:03d}번.pptx", # 언더바+3자리
        f"{num}번.pptx",
        f"{num:03d}번.pptx",
    ]:
        p = folder_path / fname
        if p.exists():
            return str(p)
    # 그래도 없으면 폴더 전체 순회해서 번호 포함 파일 탐색
    import re as _re
    for f in folder_path.iterdir():
        if f.suffix.lower() != '.pptx': continue
        if _re.search(rf'(?<!\d){num}(?!\d)', f.stem):
            return str(f)
    return None


# ───────────────────────────────────────────────
# ZIP 레벨 슬라이드 교체
# ───────────────────────────────────────────────

def _zip_delete_slides(pptx_path: str, slide_targets: list[str]):
    """ZIP 파일에서 제거된 슬라이드 XML/rels 삭제 + 중복 파일 정리."""
    import zipfile, tempfile, shutil as _sh
    from pathlib import Path as _P

    del_set = set()
    for t in slide_targets:
        t = t.lstrip('/')
        del_set.add('ppt/' + t)
        del_set.add('ppt/' + t.replace('slides/', 'slides/_rels/') + '.rels')

    with tempfile.TemporaryDirectory() as td:
        tmp = _P(td) / 'tmp.pptx'
        seen = set()
        with zipfile.ZipFile(pptx_path, 'r') as zin, \
             zipfile.ZipFile(str(tmp), 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                fn = item.filename
                if fn in del_set or fn in seen:
                    continue
                seen.add(fn)
                zout.writestr(item, zin.read(fn))
        _sh.move(str(tmp), pptx_path)



# ───────────────────────────────────────────────────────────────
# 설교찬양 (악보 이미지 PPT → 예배 자막 슬라이드 교체)
# ───────────────────────────────────────────────────────────────

def find_sermon_song_region(prs) -> list[int]:
    """
    설교찬양 슬라이드 구역 탐지.
    '직사각형 8' + '그림 1' 패턴인 슬라이드 중 첫 번째 연속 그룹 반환.
    (주제찬양은 빈 슬라이드 이후 두 번째 그룹)
    """
    region = []
    prev_idx = -2
    for i in range(len(prs.slides)):
        sl = prs.slides[i]
        names = [sh.name for sh in sl.shapes]
        if '직사각형 8' in names and '그림 1' in names:
            if region and i > prev_idx + 1:
                break   # 연속이 끊기면 첫 번째 그룹 완료
            region.append(i)
            prev_idx = i
    return region


def apply_sermon_song(prs, score_pptx_path: str, song_title: str) -> int:
    """
    설교찬양 악보 PPT의 이미지를 예배 자막에 삽입.
    - score_pptx_path: 악보 원본 PPT (.ppt/.pptx)
    - song_title: 설교찬양 제목
    반환: 생성된 슬라이드 수 (실패 시 -1 또는 0)
    """
    import warnings as _w
    import os as _os
    import tempfile as _tmp
    import shutil as _sh

    # .PPT → .pptx 변환 (필요시)
    score_path = score_pptx_path
    tmp_converted = None
    if score_pptx_path.lower().endswith('.ppt') and not score_pptx_path.lower().endswith('.pptx'):
        lo_candidates = [
            'libreoffice',
            'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
            'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe',
        ]
        converted_ok = False
        for lo_cmd in lo_candidates:
            try:
                import subprocess
                out_dir = _os.path.dirname(score_pptx_path) or _tmp.gettempdir()
                subprocess.run(
                    [lo_cmd, '--headless', '--convert-to', 'pptx',
                     score_pptx_path, '--outdir', out_dir],
                    capture_output=True, timeout=60
                )
                base = _os.path.splitext(_os.path.basename(score_pptx_path))[0]
                converted = _os.path.join(out_dir, base + '.pptx')
                if _os.path.exists(converted):
                    tmp_converted = converted
                    score_path = tmp_converted
                    converted_ok = True
                    break
            except Exception:
                continue
        if not converted_ok:
            return -1  # 변환 실패

    try:
        prs_score = Presentation(score_path)
    except Exception:
        return 0
    finally:
        if tmp_converted and _os.path.exists(tmp_converted):
            try: _os.unlink(tmp_converted)
            except: pass

    # 악보 이미지 blob 목록 추출
    score_images = []
    for sl_s in prs_score.slides:
        for sh in sl_s.shapes:
            if sh.shape_type == 13:  # PICTURE
                blip = sh._element.find('.//' + qn('a:blip'))
                if blip is not None:
                    rId = blip.get(qn('r:embed'))
                    if rId in sl_s.part.rels:
                        score_images.append(sl_s.part.rels[rId].target_part.blob)
                        break

    if not score_images:
        return 0

    # 설교찬양 구역 탐지
    region = find_sermon_song_region(prs)
    if not region:
        return 0

    # 슬라이드 수 부족하면 복제
    sldIdLst = prs.slides._sldIdLst
    while len(region) < len(score_images):
        src_sl = prs.slides[region[-1]]
        new_sl = prs.slides.add_slide(src_sl.slide_layout)
        sp = new_sl.shapes._spTree
        for c in list(sp): sp.remove(c)
        for c in src_sl.shapes._spTree: sp.append(copy.deepcopy(c))
        for rel in src_sl.part.rels.values():
            if 'slideLayout' not in rel.reltype:
                try: new_sl.part.relate_to(rel.target_part, rel.reltype)
                except: pass
        cur_ids = list(sldIdLst.sldId_lst)
        new_elem = cur_ids[-1]
        sldIdLst.remove(new_elem)
        cur_ids2 = list(sldIdLst.sldId_lst)
        last = region[-1]
        if last + 1 < len(cur_ids2):
            cur_ids2[last + 1].addprevious(new_elem)
        else:
            sldIdLst.append(new_elem)
        region.append(last + 1)

    # 슬라이드 수 초과하면 제거
    if len(score_images) < len(region):
        cur_ids = list(sldIdLst.sldId_lst)
        for elem in cur_ids[region[len(score_images)]: region[-1] + 1]:
            sldIdLst.remove(elem)
        region = region[:len(score_images)]

    # 템플릿 슬라이드에서 그림 1의 위치/크기/srcRect 정보 추출
    tmpl_sl = prs.slides[region[0]]
    img_shape_tmpl = next((sh for sh in tmpl_sl.shapes if sh.name == '그림 1'), None)
    if img_shape_tmpl is None:
        return 0

    img_left   = img_shape_tmpl.left
    img_top    = img_shape_tmpl.top
    img_width  = img_shape_tmpl.width
    img_height = img_shape_tmpl.height
    # srcRect.t 값
    blipFill_tmpl = img_shape_tmpl._element.find('.//' + qn('p:blipFill'))
    srcRect_tmpl  = blipFill_tmpl.find(qn('a:srcRect')) if blipFill_tmpl is not None else None
    src_t = srcRect_tmpl.get('t', '14180') if srcRect_tmpl is not None else '14180'

    import io as _io

    # 각 슬라이드에 악보 이미지 삽입
    for pos, img_blob in enumerate(score_images):
        sl = prs.slides[region[pos]]

        # 기존 '그림 1' shape 제거하고 새 이미지로 교체
        # add_picture로 추가 → ImagePart가 독립적으로 생성됨
        old_img_sh = next((sh for sh in sl.shapes if sh.name == '그림 1'), None)
        if old_img_sh is None:
            continue

        sp_tree = sl.shapes._spTree

        # 새 이미지 추가 (동일한 위치/크기)
        new_pic = sl.shapes.add_picture(
            _io.BytesIO(img_blob),
            left=img_left, top=img_top,
            width=img_width, height=img_height
        )

        # srcRect.t 설정 (상단 crop)
        blipFill_new = new_pic._element.find('.//' + qn('p:blipFill'))
        if blipFill_new is not None:
            srcRect_new = etree.Element(qn('a:srcRect'))
            srcRect_new.set('t', src_t)
            stretch = blipFill_new.find(qn('a:stretch'))
            if stretch is not None:
                stretch.addprevious(srcRect_new)
            else:
                blipFill_new.append(srcRect_new)

        # 기존 shape 제거, 새 shape 이름 변경
        sp_tree.remove(old_img_sh._element)
        cNvPr = new_pic._element.find('.//' + qn('p:cNvPr'))
        if cNvPr is not None:
            cNvPr.set('name', '그림 1')

        # 제목 텍스트 교체 (앞 공백 유지)
        for sh in sl.shapes:
            if sh.name == '직사각형 8' and sh.has_text_frame:
                para = sh.text_frame.paragraphs[0]
                full_text = para.text
                leading = len(full_text) - len(full_text.lstrip())
                prefix = full_text[:leading]
                if para.runs:
                    para.runs[0].text = prefix + song_title
                    for run in para.runs[1:]:
                        run.text = ''
                break

    return len(score_images)



# ───────────────────────────────────────────────
# 슬라이드 역할 탐지 & 블록 위치 계산
# ───────────────────────────────────────────────

def get_slide_role(slide) -> str:
    texts = " ".join(
        shape.text_frame.text.strip()
        for shape in slide.shapes
        if shape.has_text_frame
    )
    if not texts.strip():
        return "blank"
    if "찬송가" in texts and "장" in texts:
        return "hymn"
    if "교독문" in texts:
        return "responsory"
    if any(k in texts for k in ["나는 전능하신", "성령으로 잉태", "장사된 지", "죄를 용서받는"]):
        return "creed"
    return "other"


def find_hymn_blocks(prs) -> list[tuple[int,int]]:
    """찬송가 블록 (앞blank+가사N장+뒤blank) 위치 목록 반환 [(start_idx, count), ...]"""
    roles = [get_slide_role(prs.slides[i]) for i in range(len(prs.slides))]
    results = []
    i = 0
    while i < len(roles):
        if roles[i] == "hymn":
            start = i - 1 if i > 0 and roles[i-1] == "blank" else i
            end = i
            while end + 1 < len(roles) and roles[end+1] == "hymn":
                end += 1
            if end + 1 < len(roles) and roles[end+1] == "blank":
                end += 1
            results.append((start, end - start + 1))
            i = end + 1
        else:
            i += 1
    return results


def find_responsory_block(prs) -> tuple[int,int] | None:
    """
    교독문 전체 블록 (앞blank+제목+본문N장+뒤blank) 위치 반환 (start_idx, count).
    교독문 본문 슬라이드는 role=other로 분류되므로
    responsory 이후 creed/hymn 직전 blank까지 포함.
    """
    roles = [get_slide_role(prs.slides[i]) for i in range(len(prs.slides))]
    try:
        resp_idx = next(i for i, r in enumerate(roles) if r == "responsory")
    except StopIteration:
        return None
    # 앞 blank 포함
    start = resp_idx - 1 if resp_idx > 0 and roles[resp_idx-1] == "blank" else resp_idx
    # 뒤: creed/hymn 나오기 직전까지, 첫 blank에서 멈춤
    end = resp_idx
    while end + 1 < len(roles):
        next_role = roles[end + 1]
        if next_role in ("creed", "hymn"):
            break
        end += 1
        if next_role == "blank":
            break
    return (start, end - start + 1)


# ───────────────────────────────────────────────
# 텍스트 치환 유틸
# ───────────────────────────────────────────────

def _replace_date_in_prs(prs, date_str: str):
    """
    전체 슬라이드에서 날짜 치환.
    날짜가 run 여러 개로 쪼개진 경우도 처리:
    예) run[0]='2026' run[1]='년 ' run[2]='1' run[3]='월 ' run[4]='4' run[5]='일'
    → paragraph 전체 텍스트를 합쳐서 날짜 패턴 탐지 후 첫 run에 전체 날짜 삽입,
      나머지 날짜 관련 run은 비움.
    """
    DATE_PAT = re.compile(r'\d{4}년\s*\d{1,2}월\s*\d{1,2}일')
    for i in range(len(prs.slides)):
        sl = prs.slides[i]
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                runs = para.runs
                if not runs:
                    continue
                # 단일 run에 날짜 전체가 있는 경우
                for run in runs:
                    if DATE_PAT.search(run.text):
                        run.text = DATE_PAT.sub(date_str, run.text)
                        break
                else:
                    # run이 쪼개진 경우: para 전체 텍스트 합쳐서 확인
                    full = ''.join(r.text for r in runs)
                    if not DATE_PAT.search(full):
                        continue
                    # 날짜 패턴에 해당하는 run 범위 찾기
                    # 연도 숫자로 시작하는 run 탐색
                    year_pat = re.compile(r'^\d{4}$')
                    start_idx = next(
                        (k for k, r in enumerate(runs) if year_pat.match(r.text.strip())),
                        None
                    )
                    if start_idx is None:
                        continue
                    # '일' 로 끝나는 run까지 범위 탐색
                    end_idx = start_idx
                    for k in range(start_idx, min(start_idx + 8, len(runs))):
                        if '일' in runs[k].text:
                            end_idx = k
                            break
                    # 첫 run에 전체 날짜 삽입, 나머지 비움
                    runs[start_idx].text = date_str
                    for k in range(start_idx + 1, end_idx + 1):
                        runs[k].text = ''


def replace_text_all(prs, old: str, new: str):
    """전체 슬라이드에서 텍스트 치환"""
    for i in range(len(prs.slides)):
        slide = prs.slides[i]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)



def _is_sermon_title_shape(tf) -> bool:
    """
    설교 제목 shape 판별.
    조건: ( ) 포함 + 최대 폰트가 42~50pt 범위 (교독문=60pt, 축도=72pt 제외)
    """
    for para in tf.paragraphs:
        runs = para.runs
        if len(runs) < 3:
            continue
        if not any('(' in r.text for r in runs):
            continue
        sizes = []
        for run in runs:
            rpr = run._r.find(qn('a:rPr'))
            if rpr is not None and rpr.get('sz'):
                sizes.append(int(rpr.get('sz')) // 100)
        if not sizes:
            continue
        max_sz = max(sizes)
        # 설교제목: 최대 폰트 42~52pt, 두 가지 크기 존재
        if 42 <= max_sz <= 52 and len(set(sizes)) >= 2:
            return True
    return False


def replace_sermon_title(prs, sermon_title: str, sermon_ref: str) -> bool:
    for i in range(len(prs.slides)):
        sl = prs.slides[i]
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            if not _is_sermon_title_shape(shape.text_frame):
                continue
            para = shape.text_frame.paragraphs[0]
            p_elem = para._p
            runs = para.runs
            if not runs:
                continue
            title_rpr = ref_rpr = None
            found_paren = False
            for run in runs:
                rpr = run._r.find(qn('a:rPr'))
                if '(' in run.text:
                    found_paren = True
                if not found_paren and rpr is not None and title_rpr is None:
                    title_rpr = copy.deepcopy(rpr)
                if found_paren and rpr is not None and ref_rpr is None:
                    ref_rpr = copy.deepcopy(rpr)
            if title_rpr is None or ref_rpr is None:
                continue
            for r in p_elem.findall(qn('a:r')):
                p_elem.remove(r)
            def _mkrun(text, rpr_tmpl):
                r = etree.Element(qn('a:r'))
                r.append(copy.deepcopy(rpr_tmpl))
                t_elem = etree.SubElement(r, qn('a:t'))
                t_elem.text = text
                return r
            ref_text = sermon_ref if sermon_ref else ''
            new_runs = [
                _mkrun(sermon_title + ' ', title_rpr),
                _mkrun('(', ref_rpr),
                _mkrun(ref_text, ref_rpr),
                _mkrun(')', ref_rpr),
            ]
            end_rpr = p_elem.find(qn('a:endParaRPr'))
            if end_rpr is not None:
                idx = list(p_elem).index(end_rpr)
                for k, nr in enumerate(new_runs):
                    p_elem.insert(idx + k, nr)
            else:
                for nr in new_runs:
                    p_elem.append(nr)
            return True
    return False

def split_into_blocks(text: str) -> list[list[str]]:
    """빈 줄로 구분된 텍스트를 슬라이드 블록 목록으로 변환"""
    blocks = []
    current = []
    for line in text.splitlines():
        if line.strip():
            current.append(line.strip())
        else:
            if current:
                blocks.append(current)
                current = []
    if current:
        blocks.append(current)
    return blocks

# ───────────────────────────────────────────────
# ZIP 슬라이드 교체 (찬송가/교독문 블록 단위 교체)
# ───────────────────────────────────────────────

def replace_slides_zip(dst_path: str, start_idx: int, remove_count: int,
                       src_path: str, out_path: str) -> int:
    """
    dst_path의 start_idx ~ start_idx+remove_count-1 슬라이드를
    src_path 전체 슬라이드로 교체 후 out_path에 저장.
    반환: 새로 삽입된 슬라이드 수
    """
    import zipfile, re as _re, tempfile, shutil
    from pathlib import Path

    with tempfile.TemporaryDirectory() as tmpdir:
        dst_dir = Path(tmpdir) / 'dst'
        src_dir = Path(tmpdir) / 'src'
        with zipfile.ZipFile(dst_path, 'r') as z: z.extractall(dst_dir)
        with zipfile.ZipFile(src_path, 'r') as z: z.extractall(src_dir)

        dst_prs_xml  = dst_dir / 'ppt' / 'presentation.xml'
        dst_rels_xml = dst_dir / 'ppt' / '_rels' / 'presentation.xml.rels'
        dst_ct_xml   = dst_dir / '[Content_Types].xml'

        dst_prs_root  = etree.parse(str(dst_prs_xml)).getroot()
        dst_rels_root = etree.parse(str(dst_rels_xml)).getroot()
        dst_ct_root   = etree.parse(str(dst_ct_xml)).getroot()

        sldIdLst = dst_prs_root.find(f'{{{NS_P}}}sldIdLst')
        rids = [int(m.group(1)) for e in dst_rels_root.findall(f'{{{NS_REL}}}Relationship')
                if (m := _re.match(r'rId(\d+)', e.get('Id', '')))]
        sids = [int(e.get('id', 0)) for e in sldIdLst]
        next_rid = max(rids, default=0) + 1
        next_sid = max(sids, default=255) + 1

        src_prs_root  = etree.parse(str(src_dir / 'ppt' / 'presentation.xml')).getroot()
        src_rels_root = etree.parse(str(src_dir / 'ppt' / '_rels' / 'presentation.xml.rels')).getroot()
        src_sldIdLst  = src_prs_root.find(f'{{{NS_P}}}sldIdLst')

        src_rid_to_target = {
            e.get('Id'): e.get('Target', '')
            for e in src_rels_root.findall(f'{{{NS_REL}}}Relationship')
            if 'slide' in e.get('Type', '') and 'slideLayout' not in e.get('Type', '')
        }
        src_slide_files = [
            Path(src_rid_to_target.get(sld.get(f'{{{NS_R}}}id'), '')).name
            for sld in list(src_sldIdLst)
            if Path(src_rid_to_target.get(sld.get(f'{{{NS_R}}}id'), '')).name
        ]

        dst_slides_dir = dst_dir / 'ppt' / 'slides'
        dst_media_dir  = dst_dir / 'ppt' / 'media'
        dst_media_dir.mkdir(exist_ok=True)
        existing_media = {f.name for f in dst_media_dir.iterdir()}

        # 제거 대상 파일명 수집
        rid_to_target = {e.get('Id'): e.get('Target', '')
                         for e in dst_rels_root.findall(f'{{{NS_REL}}}Relationship')}
        files_to_remove = []
        for elem in list(sldIdLst)[start_idx: start_idx + remove_count]:
            rid = elem.get(f'{{{NS_R}}}id')
            target = rid_to_target.get(rid, '')
            if target:
                files_to_remove.append(Path(target).name)

        new_sldId_elems = []

        for fname in src_slide_files:
            new_fname = f'slide_new_{next_rid}.xml'
            dst_slide_path = dst_slides_dir / new_fname
            dst_rels_dir2  = dst_slides_dir / '_rels'
            dst_rels_dir2.mkdir(exist_ok=True)
            shutil.copy2(src_dir / 'ppt' / 'slides' / fname, dst_slide_path)

            src_rels_p = src_dir / 'ppt' / 'slides' / '_rels' / (fname + '.rels')
            if src_rels_p.exists():
                rr = etree.parse(str(src_rels_p)).getroot()
                for rel in rr.findall(f'{{{NS_REL}}}Relationship'):
                    t = rel.get('Target', '')
                    if '../media/' in t:
                        mname = Path(t).name
                        src_media = src_dir / 'ppt' / 'media' / mname
                        if src_media.exists():
                            dst_mname = mname if mname not in existing_media else f'src_{next_rid}_{mname}'
                            shutil.copy2(src_media, dst_media_dir / dst_mname)
                            existing_media.add(dst_mname)
                            rel.set('Target', f'../media/{dst_mname}')
                with open(dst_rels_dir2 / (new_fname + '.rels'), 'wb') as f:
                    f.write(etree.tostring(rr, xml_declaration=True, encoding='UTF-8', standalone=True))

            new_rid_str = f'rId{next_rid}'
            rel_e = etree.SubElement(dst_rels_root, f'{{{NS_REL}}}Relationship')
            rel_e.set('Id', new_rid_str)
            rel_e.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
            rel_e.set('Target', f'slides/{new_fname}')

            pn = f'/ppt/slides/{new_fname}'
            if pn not in {e.get('PartName') for e in dst_ct_root.findall(f'{{{NS_CT}}}Override')}:
                ov = etree.SubElement(dst_ct_root, f'{{{NS_CT}}}Override')
                ov.set('PartName', pn)
                ov.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')

            new_sid_elem = etree.Element(f'{{{NS_P}}}sldId')
            new_sid_elem.set('id', str(next_sid))
            new_sid_elem.set(f'{{{NS_R}}}id', new_rid_str)
            new_sldId_elems.append(new_sid_elem)
            next_rid += 1
            next_sid += 1

        # 기존 슬라이드 제거 후 새 슬라이드 삽입
        all_ids = list(sldIdLst)
        for elem in all_ids[start_idx: start_idx + remove_count]:
            sldIdLst.remove(elem)
        all_ids_after = list(sldIdLst)
        if start_idx < len(all_ids_after):
            ref = all_ids_after[start_idx]
            for ne in new_sldId_elems:
                ref.addprevious(ne)
        else:
            for ne in new_sldId_elems:
                sldIdLst.append(ne)

        # 제거된 슬라이드 파일 삭제
        for fname in files_to_remove:
            for p in [dst_slides_dir / fname,
                      dst_slides_dir / '_rels' / (fname + '.rels')]:
                if p.exists(): p.unlink()

        for path, root in [(dst_prs_xml,  dst_prs_root),
                           (dst_rels_xml, dst_rels_root),
                           (dst_ct_xml,   dst_ct_root)]:
            with open(path, 'wb') as f:
                f.write(etree.tostring(root, xml_declaration=True,
                                       encoding='UTF-8', standalone=True))

        with zipfile.ZipFile(out_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for file in sorted(dst_dir.rglob('*')):
                if file.is_file():
                    zout.write(file, file.relative_to(dst_dir))

    return len(src_slide_files)


# ───────────────────────────────────────────────
# 텍스트박스 유틸
# ───────────────────────────────────────────────

def _set_choir_cover_text(shape, choir_name: str, song_title: str):
    """
    성가대 표지 shape 텍스트 교체.
    para[0]: 성가대 이름, para[1]: 빈 줄, para[2]: 곡 제목
    """
    tf = shape.text_frame
    paras = tf.paragraphs
    if len(paras) < 3:
        _set_textbox_lines(shape, [choir_name, '', song_title])
        return
    p0_runs = paras[0].runs
    if p0_runs:
        p0_runs[0].text = choir_name
        for extra_run in p0_runs[1:]:
            extra_run.text = ''
    p2_runs = paras[2].runs
    if p2_runs:
        p2_runs[0].text = song_title
    else:
        p2_elem = paras[2]._p
        rpr_tmpl = paras[0].runs[0]._r.find(qn('a:rPr')) if paras[0].runs else None
        r_elem = etree.SubElement(p2_elem, qn('a:r'))
        if rpr_tmpl is not None:
            r_elem.insert(0, copy.deepcopy(rpr_tmpl))
        t_elem = etree.SubElement(r_elem, qn('a:t'))
        t_elem.text = song_title


def _set_textbox_lines(shape, lines: list[str], anchor: str = None):
    """TextBox의 텍스트를 줄 목록으로 교체 (기존 서식 유지).
    anchor: 't'=상단, 'ctr'=가운데, None=변경 없음
    """
    tf = shape.text_frame
    txBody = tf._txBody
    existing_paras = txBody.findall(qn('a:p'))
    if not existing_paras:
        return
    tmpl_p = existing_paras[0]
    tmpl_runs = tmpl_p.findall(qn('a:r'))
    tmpl_rpr = copy.deepcopy(tmpl_runs[0].find(qn('a:rPr'))) if tmpl_runs else None
    tmpl_ppr = copy.deepcopy(tmpl_p.find(qn('a:pPr')))
    end_rpr = copy.deepcopy(tmpl_p.find(qn('a:endParaRPr')))

    for p in existing_paras:
        txBody.remove(p)

    for line in lines:
        p_e = etree.Element(qn('a:p'))
        if tmpl_ppr is not None:
            p_e.append(copy.deepcopy(tmpl_ppr))
        if line:
            r_e = etree.SubElement(p_e, qn('a:r'))
            if tmpl_rpr is not None:
                r_e.insert(0, copy.deepcopy(tmpl_rpr))
            t_e = etree.SubElement(r_e, qn('a:t'))
            t_e.text = line
        if end_rpr is not None:
            p_e.append(copy.deepcopy(end_rpr))
        txBody.append(p_e)
    if anchor is not None:
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', anchor)


def _split_title(title: str) -> list[str]:
    """찬양 제목이 길면 두 줄로 분할. [] 안에 줄바꿈이 있으면 그대로 사용."""
    if '\n' in title:
        parts = [p.strip() for p in title.split('\n') if p.strip()]
        return parts[:2] if parts else [title]
    if len(title) <= 7:
        return [title]
    mid = len(title) // 2
    sp = title.rfind(' ', 0, mid + 3)
    return [title[:sp], title[sp+1:]] if sp > 0 else [title]


# ───────────────────────────────────────────────
# 파싱 함수
# ───────────────────────────────────────────────

def parse_praise_input(text: str) -> list[dict]:
    """
    [제목] 구분으로 찬양팀 가사 파싱.
    [] 안에 줄바꿈이 있어도 제목으로 인식하며 줄바꿈 유지.
    반환: [{'title': str, 'slides': [['줄1','줄2'], ...]}, ...]
    """
    import re as _re
    # 웹 입력의 \r\n 또는 \r 줄바꿈을 \n으로 통일
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    NL = '\x00'
    def _fix_bracket(m):
        inner = m.group(1).replace('\n', NL).strip(NL).strip()
        return '[' + inner + ']'
    text = _re.sub(r'\[([^\]]*?)\]', _fix_bracket, text, flags=_re.DOTALL)

    songs = []
    current_title = ""
    current_slides = []
    current_lines = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if line.startswith('[') and line.endswith(']'):
            if current_lines:
                current_slides.append(current_lines)
                current_lines = []
            if current_title or current_slides:
                songs.append({'title': current_title, 'slides': current_slides})
            current_title = line[1:-1].replace(NL, '\n').strip()
            current_slides = []
        elif line == '':
            if current_lines:
                current_slides.append(current_lines)
                current_lines = []
        else:
            current_lines.append(line)
    if current_lines:
        current_slides.append(current_lines)
    if current_title or current_slides:
        songs.append({'title': current_title, 'slides': current_slides})
    return songs


def parse_choir_input(text: str) -> list:
    """
    빈 줄로 슬라이드 구분하여 성가대 가사 파싱.
    (간주) 등 괄호 단어는 빈 슬라이드 마커로 처리.
    반환: list of (list[str] | 'interlude')
    """
    import re as _re
    INTERLUDE_PAT = _re.compile(r'^\(간주\)$')
    blocks = []
    current = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if INTERLUDE_PAT.match(line):
            if current:
                blocks.append(current)
                current = []
            blocks.append('interlude')
        elif line == '':
            if current:
                blocks.append(current)
                current = []
        else:
            current.append(line)
    if current:
        blocks.append(current)
    return blocks


def parse_godpia_text(text: str, book_ref: str) -> list[dict]:
    """
    Godpia 복사 텍스트 파싱.
    형식: '1본문내용\n2본문내용...'
    book_ref: '룻기 3' 형태 (책이름+장)
    반환: [{'label': '룻기 3:1', 'lines': ['전체구절텍스트']}, ...]
    줄 분리 없이 전체를 한 줄로 넘김 — PPT 텍스트박스 폭에 맞게 자동 줄바꿈됨
    """
    import re as _re
    results = []
    parts = _re.split(r'(?m)^(\d+)(?=[가-힣])', text.strip())
    i = 1
    while i < len(parts) - 1:
        verse_num = parts[i].strip()
        verse_text = parts[i + 1].strip()
        label = f"{book_ref}:{verse_num}"
        results.append({'label': label, 'lines': [verse_text]})
        i += 2
    return results


# ───────────────────────────────────────────────
# 구역 탐지
# ───────────────────────────────────────────────

def get_bible_region(prs) -> tuple[int, int] | tuple[None, None]:
    """성경봉독 구역 (직사각형 3 + 직사각형 4) 슬라이드 범위 반환."""
    import re as _re
    REF_PAT = _re.compile(r'[가-힣]+\s*\d+:\d+')
    for i, sl in enumerate(prs.slides):
        names = [sh.name for sh in sl.shapes]
        if '직사각형 3' in names and '직사각형 4' in names:
            for sh in sl.shapes:
                if sh.name == '직사각형 4' and sh.has_text_frame:
                    if REF_PAT.search(sh.text_frame.text.strip()):
                        start = i
                        end = i
                        while end + 1 < len(prs.slides):
                            next_names = [s.name for s in prs.slides[end+1].shapes]
                            if '직사각형 3' in next_names and '직사각형 4' in next_names:
                                end += 1
                            else:
                                break
                        return start, end
    return None, None


def apply_bible_text(prs, verses: list[dict]) -> int:
    """
    성경봉독 슬라이드에 구절 적용.
    verses: [{'label': '룻기 3:1', 'lines': ['줄1', '줄2']}, ...]
    반환: 적용된 슬라이드 수
    """
    start, end = get_bible_region(prs)
    if start is None or not verses:
        return 0

    region = list(range(start, end + 1))
    needed = len(verses)
    sldIdLst = prs.slides._sldIdLst

    # 슬라이드 수 조정
    while len(region) < needed:
        last = region[-1]
        src_sl = prs.slides[last]
        new_sl = prs.slides.add_slide(src_sl.slide_layout)
        sp = new_sl.shapes._spTree
        for c in list(sp): sp.remove(c)
        for c in src_sl.shapes._spTree: sp.append(copy.deepcopy(c))
        for rel in src_sl.part.rels.values():
            if 'slideLayout' not in rel.reltype:
                try: new_sl.part.relate_to(rel.target_part, rel.reltype)
                except: pass
        cur_ids = list(sldIdLst.sldId_lst)
        new_elem = cur_ids[-1]
        sldIdLst.remove(new_elem)
        cur_ids2 = list(sldIdLst.sldId_lst)
        if last + 1 < len(cur_ids2):
            cur_ids2[last + 1].addprevious(new_elem)
        else:
            sldIdLst.append(new_elem)
        region.append(last + 1)

    if needed < len(region):
        cur_ids = list(sldIdLst.sldId_lst)
        for elem in cur_ids[start + needed: start + len(region)]:
            sldIdLst.remove(elem)
        region = region[:needed]

    for pos, verse in enumerate(verses):
        if pos >= len(region):
            break
        sl = prs.slides[region[pos]]
        lines = verse['lines']
        label = verse['label']
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            if sh.name == '직사각형 3':
                _set_textbox_lines(sh, lines)  # 줄바꿈은 PPT 텍스트박스에 맡김
            elif sh.name == '직사각형 4':
                _set_textbox_lines(sh, [label])

    return len(verses)


# ───────────────────────────────────────────────
# 찬양팀 (2부)
# ───────────────────────────────────────────────

def get_praise_region(prs) -> tuple[int, int] | tuple[None, None]:
    """찬양팀 구역 (자유형: 도형 14) 슬라이드 범위 반환."""
    for i, sl in enumerate(prs.slides):
        if any(sh.name == '자유형: 도형 14' for sh in sl.shapes):
            start = i
            end = i
            # 그룹4(구분자) 슬라이드까지 포함한 연속 구역
            while end + 1 < len(prs.slides):
                next_sl = prs.slides[end + 1]
                next_names = [sh.name for sh in next_sl.shapes]
                if '그룹 4' in next_names or '자유형: 도형 14' in next_names:
                    end += 1
                else:
                    break
            # 앞 구분자(그룹4) 포함
            if start > 0:
                prev_names = [sh.name for sh in prs.slides[start-1].shapes]
                if '그룹 4' in prev_names:
                    start -= 1
            return start, end
    return None, None


def _build_praise_src_pptx(template_path: str, songs: list[dict]) -> str:
    """
    찬양팀 가사를 담은 임시 pptx 생성 후 경로 반환.
    replace_slides_zip의 src 파일로 사용.
    """
    import tempfile as _tmp, shutil as _sh2, os as _os2

    NS_P_local = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    WIPE = (f'<p:transition xmlns:p="{NS_P_local}" spd="slow">'
            '<p:wipe dir="r"/></p:transition>')

    prs = Presentation(template_path)
    start, end = get_praise_region(prs)
    if start is None:
        return None

    sep_slide   = prs.slides[start]
    lyric_slide = next(
        (prs.slides[i] for i in range(start, end + 1)
         if any(sh.name == '자유형: 도형 14' for sh in prs.slides[i].shapes)), None)
    if lyric_slide is None:
        return None

    sep_spTree   = copy.deepcopy(sep_slide.shapes._spTree)
    lyric_spTree = copy.deepcopy(lyric_slide.shapes._spTree)
    sep_rels   = [(rel.reltype, rel.target_part)
                  for rel in sep_slide.part.rels.values()
                  if any(t in rel.reltype for t in ('image', 'hdphoto'))]
    lyric_rels = [(rel.reltype, rel.target_part)
                  for rel in lyric_slide.part.rels.values()
                  if any(t in rel.reltype for t in ('image', 'hdphoto'))]

    def _get_fmt(spTree, name):
        for sp in spTree.iter(qn('p:sp')):
            cnv = sp.find('.//' + qn('p:cNvPr'))
            if cnv is not None and cnv.get('name') == name:
                paras = sp.findall('.//' + qn('a:p'))
                if paras:
                    runs = paras[0].findall(qn('a:r'))
                    if runs:
                        rpr = runs[0].find(qn('a:rPr'))
                        ppr = paras[0].find(qn('a:pPr'))
                        return (copy.deepcopy(rpr) if rpr is not None else None,
                                copy.deepcopy(ppr) if ppr is not None else None)
        return (None, None)

    rpr14, ppr14 = _get_fmt(lyric_spTree, '자유형: 도형 14')
    rpr6,  ppr6  = _get_fmt(lyric_spTree, '자유형: 도형 6')

    # 원본 슬라이드 레이아웃 인덱스
    sep_layout_idx   = next((i for i, l in enumerate(prs.slide_layouts)
                              if l == sep_slide.slide_layout), 0)
    lyric_layout_idx = next((i for i, l in enumerate(prs.slide_layouts)
                              if l == lyric_slide.slide_layout), 6)

    # 임시 pptx — 원본 파일 복사해서 슬라이드 비우기
    blank_path = _tmp.mktemp(suffix='.pptx')
    _sh2.copy2(template_path, blank_path)
    blank_prs = Presentation(blank_path)
    sldIdLst_b = blank_prs.slides._sldIdLst
    for elem in list(sldIdLst_b.sldId_lst):
        sldIdLst_b.remove(elem)
    layout_sep   = blank_prs.slide_layouts[sep_layout_idx]   if sep_layout_idx < len(blank_prs.slide_layouts) else blank_prs.slide_layouts[0]
    layout_lyric = blank_prs.slide_layouts[lyric_layout_idx] if lyric_layout_idx < len(blank_prs.slide_layouts) else blank_prs.slide_layouts[0]

    # 원본 전환효과 XML 추출
    _mc_ns = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
    _sep_alt   = sep_slide._element.find(f'{{{_mc_ns}}}AlternateContent')
    sep_alt_xml = etree.tostring(_sep_alt).decode() if _sep_alt is not None else None
    _lyric2_idx = next((i for i in range(start, end + 1)
                        if i != next((j for j in range(start, end + 1)
                                      if any(sh.name == '자유형: 도형 14'
                                             for sh in prs.slides[j].shapes)), start)
                        and any(sh.name == '자유형: 도형 14' for sh in prs.slides[i].shapes)), None)
    _lyric2_alt = (prs.slides[_lyric2_idx]._element.find(f'{{{_mc_ns}}}AlternateContent')
                   if _lyric2_idx is not None else None)
    lyric2_alt_xml = etree.tostring(_lyric2_alt).decode() if _lyric2_alt is not None else None

    def _add(spTree_tmpl, rels, with_wipe, is_sep=False):
        sl = blank_prs.slides.add_slide(layout_sep if is_sep else layout_lyric)
        sp = sl.shapes._spTree
        for c in list(sp): sp.remove(c)
        for c in spTree_tmpl: sp.append(copy.deepcopy(c))
        for reltype, tpart in rels:
            try: sl.part.relate_to(tpart, reltype)
            except: pass
        if with_wipe:
            sl._element.append(etree.fromstring(WIPE))
        elif is_sep and sep_alt_xml:
            sl._element.append(etree.fromstring(sep_alt_xml))
        elif not is_sep and lyric2_alt_xml:
            sl._element.append(etree.fromstring(lyric2_alt_xml))
        return sl

    def _set_anchor(shape, anchor):
        """텍스트박스 세로 정렬 설정 (t=상단, ctr=가운데, b=하단)."""
        bodyPr = shape.text_frame._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', anchor)

    def _set_text(sl, lines, title):
        for sh in sl.shapes:
            if not sh.has_text_frame: continue
            if sh.name == '자유형: 도형 14':
                # 1줄이면 빈줄 추가해서 항상 2줄 유지 (가사 위쪽 정렬 효과)
                padded = lines[:2] if len(lines) >= 2 else [lines[0], ''] if lines else ['', '']
                _raw(sh, padded, rpr14, ppr14)
            elif sh.name == '자유형: 도형 6':
                title_lines = _split_title(title)
                # 한 줄이 6자 초과이면 자간 좁게(spc=-100) 적용
                if any(len(l) > 6 for l in title_lines):
                    rpr6_use = copy.deepcopy(rpr6) if rpr6 is not None else etree.Element(qn('a:rPr'))
                    rpr6_use.set('spc', '-100')
                    _raw(sh, title_lines, rpr6_use, ppr6)
                else:
                    _raw(sh, title_lines, rpr6, ppr6)

    def _raw(shape, lines, rpr_t, ppr_t):
        txBody = shape.text_frame._txBody
        for p in txBody.findall(qn('a:p')): txBody.remove(p)
        for line in lines:
            # XML 비호환 제어 문자 제거 (\t, \n 제외)
            clean = ''.join(c for c in (line or '') if c >= ' ' or c in '\t\n')
            p_e = etree.Element(qn('a:p'))
            if ppr_t is not None: p_e.append(copy.deepcopy(ppr_t))
            r_e = etree.SubElement(p_e, qn('a:r'))
            if rpr_t is not None: r_e.insert(0, copy.deepcopy(rpr_t))
            t_e = etree.SubElement(r_e, qn('a:t'))
            t_e.text = clean
            txBody.append(p_e)

    for song in songs:
        _add(sep_spTree, sep_rels, False, is_sep=True)
        for idx, lines in enumerate(song['slides']):
            sl = _add(lyric_spTree, lyric_rels, idx == 0, is_sep=False)
            _set_text(sl, lines, song['title'])
    _add(sep_spTree, sep_rels, False, is_sep=True)

    tmp = _tmp.mktemp(suffix='.pptx')
    blank_prs.save(tmp)
    try: _os2.unlink(blank_path)
    except: pass
    # save() 시 원본 슬라이드와 새 슬라이드가 중복 저장됨 → 중복 제거
    _zip_delete_slides(tmp, [])
    return tmp


def apply_praise_lyrics(out_path: str, songs: list[dict]) -> int:
    """
    2부 찬양팀 가사 적용. replace_slides_zip 방식으로 교체.
    out_path: 저장할 파일 경로 (in-place 수정)
    반환: 생성된 가사 슬라이드 수
    """
    import warnings as _w, os as _os

    prs_check = Presentation(out_path)
    start, end = get_praise_region(prs_check)
    if start is None or not songs:
        return 0

    src_path = _build_praise_src_pptx(out_path, songs)
    if src_path is None:
        return 0

    try:
        with _w.catch_warnings():
            _w.simplefilter("ignore")
            replace_slides_zip(out_path, start, end - start + 1, src_path, out_path)
    finally:
        try: _os.unlink(src_path)
        except: pass

    return sum(len(s['slides']) for s in songs)


# ───────────────────────────────────────────────
# 성가대
# ───────────────────────────────────────────────

def get_choir_region(prs) -> tuple[int, int] | tuple[None, None]:
    """성가대 가사 구역 (TextBox 1 + 그룹 5) 슬라이드 범위 반환.
    뒤쪽 구분자(그룹 5만 있는 빈 슬라이드)는 구역에서 제외."""
    for i, sl in enumerate(prs.slides):
        names = [sh.name for sh in sl.shapes]
        if 'TextBox 1' in names and '그룹 5' in names:
            start = i
            end = i
            # TextBox 1이 있거나 그룹5만 있는 슬라이드까지 포함
            while end + 1 < len(prs.slides):
                next_names = [sh.name for sh in prs.slides[end+1].shapes]
                if '그룹 5' in next_names:
                    end += 1
                else:
                    break
            # 뒤쪽 그룹5만 있는 슬라이드(구분자) 제외
            while end > start:
                last_names = [sh.name for sh in prs.slides[end].shapes]
                if last_names == ['그룹 5']:
                    end -= 1
                else:
                    break
            return start, end
    return None, None


def apply_choir_lyrics(prs, blocks: list[list[str]],
                        choir_name: str, song_title: str) -> int:
    """
    성가대 가사를 슬라이드에 적용.
    - blocks: 슬라이드별 가사 줄 목록. 'interlude'면 간주 빈 슬라이드.
    반환: 실제 적용된 슬라이드 수
    """
    start, end = get_choir_region(prs)
    if start is None:
        return 0

    # 성가대 표지 수정
    for cover_i in range(max(0, start - 3), start):
        sl = prs.slides[cover_i]
        for shape in sl.shapes:
            if (shape.has_text_frame and '직사각형 6' in shape.name
                    and '찬양대' in shape.text_frame.text):
                _set_choir_cover_text(shape, choir_name, song_title)
                break

    FADE_XML = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:fade/></p:transition>'

    def _add_fade(slide_elem):
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        for t in slide_elem.findall(f'{{{ns_p}}}transition'):
            slide_elem.remove(t)
        slide_elem.append(etree.fromstring(FADE_XML))

    def _get_blank_template_idx(prs_obj, region_start: int) -> int:
        for k in range(region_start - 1, max(0, region_start - 4), -1):
            sl = prs_obj.slides[k]
            names = [sh.name for sh in sl.shapes]
            if names == ['그룹 5']:
                return k
        return region_start

    def _clone_lyrics_slide(prs_obj, template_idx: int, after_idx: int):
        src = prs_obj.slides[template_idx]
        src_layout = src.slide_layout
        new_slide = prs_obj.slides.add_slide(src_layout)
        dst = new_slide.shapes._spTree
        src_tree = src.shapes._spTree
        for child in list(dst): dst.remove(child)
        for child in src_tree: dst.append(copy.deepcopy(child))
        for rel in src.part.rels.values():
            if 'image' in rel.reltype:
                try: new_slide.part.relate_to(rel.target_part, rel.reltype)
                except: pass
        sldIdLst = prs_obj.slides._sldIdLst
        all_ids = list(sldIdLst.sldId_lst)
        new_elem = all_ids[-1]
        sldIdLst.remove(new_elem)
        cur_ids = list(sldIdLst.sldId_lst)
        if after_idx + 1 < len(cur_ids):
            cur_ids[after_idx + 1].addprevious(new_elem)
        else:
            sldIdLst.append(new_elem)
        return after_idx + 1

    region = list(range(start, end + 1))
    needed = len(blocks)
    template_idx = region[0] if region else start
    blank_template_idx = _get_blank_template_idx(prs, start)

    while len(region) < needed:
        last = region[-1]
        pos_for_this = len(region)
        is_interlude_slot = (pos_for_this < len(blocks) and blocks[pos_for_this] == 'interlude')
        tmpl = blank_template_idx if is_interlude_slot else template_idx
        new_idx = _clone_lyrics_slide(prs, tmpl, last)
        region.append(new_idx)

    applied = 0
    need_fade = False

    for pos, block in enumerate(blocks):
        if pos >= len(region):
            break
        target_idx = region[pos]
        sl = prs.slides[target_idx]

        if block == 'interlude':
            spTree = sl.shapes._spTree
            for sh in list(sl.shapes):
                if sh.name == 'TextBox 1':
                    spTree.remove(sh._element)
            # 간주 빈 슬라이드에도 fade 전환효과 추가
            _add_fade(sl._element)
            need_fade = True
            continue

        for sh in sl.shapes:
            if sh.has_text_frame and sh.name == 'TextBox 1':
                lines = block[:2]
                if len(lines) == 1:
                    lines = [lines[0], '']
                _set_textbox_lines(sh, lines)
                if need_fade:
                    _add_fade(sl._element)
                    need_fade = False
                applied += 1
                break

    if needed < len(region):
        sldIdLst = prs.slides._sldIdLst
        all_ids = list(sldIdLst.sldId_lst)
        for elem in all_ids[start + needed: start + len(region)]:
            sldIdLst.remove(elem)

    return applied



